#!/usr/bin/env python3
"""
ingest.py — 미처리 raw/ 파일을 탐지하고 상태를 관리한다.
실제 wiki 컴파일은 LLM 엔진 (claude CLI 또는 API)이 담당한다.

사용법:
  python ingest.py                          # 미처리 파일 목록 출력
  python ingest.py --url https://...        # URL 스크랩 → raw/clippings/ 저장
  python ingest.py --file ~/paper.pdf       # 로컬 파일 → raw/docs/ 저장
  python ingest.py --note "텍스트"          # 텍스트 → raw/notes/ 저장
  python ingest.py --mark-done              # 현재 raw/ 전체를 처리 완료로 표시
  python ingest.py --url ... --resonance high   # resonance 레벨 지정 저장
  python ingest.py --priority-only          # resonance: high 파일만 목록 출력
"""
import argparse
import json
import re
import shutil
import sys
from datetime import datetime
from pathlib import Path

import httpx
from markdownify import markdownify

WIKI_ROOT = Path(__file__).parent.parent
RAW_DIR = WIKI_ROOT / "raw"
STATE_FILE = WIKI_ROOT / ".ingest_state.json"

SUPPORTED_EXTENSIONS = {".md", ".txt", ".pdf", ".docx", ".pptx"}


def load_state() -> dict:
    if STATE_FILE.exists():
        return json.loads(STATE_FILE.read_text())
    return {"processed": []}


def save_state(state: dict) -> None:
    STATE_FILE.write_text(json.dumps(state, indent=2, default=str))


def extract_text(file: Path) -> str | None:
    """지원 형식에서 텍스트를 추출해 MD 문자열로 반환한다."""
    suffix = file.suffix.lower()

    if suffix in {".md", ".txt"}:
        return file.read_text(errors="replace")

    if suffix == ".pdf":
        import fitz  # pymupdf
        doc = fitz.open(str(file))
        pages = [page.get_text() for page in doc]
        doc.close()
        return "\n\n".join(pages)

    if suffix == ".docx":
        from docx import Document
        doc = Document(str(file))
        return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())

    if suffix == ".pptx":
        from pptx import Presentation
        prs = Presentation(str(file))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text_frame.text
                for shape in slide.shapes
                if shape.has_text_frame and shape.text_frame.text.strip()
            ]
            if texts:
                slides.append(f"## 슬라이드 {i}\n\n" + "\n\n".join(texts))
        return "\n\n".join(slides)

    return None


def _get_resonance(file: Path) -> str | None:
    """파일 frontmatter에서 resonance 값을 읽어 반환한다. 없으면 None."""
    if file.suffix.lower() not in {".md", ".txt"}:
        return None
    try:
        content = file.read_text(errors="replace")
        m = re.search(r"^resonance:\s*(\S+)", content, re.MULTILINE)
        return m.group(1).lower() if m else None
    except OSError:
        return None


def _merge_resonance_frontmatter(content: str, resonance: str) -> str:
    """md/txt 본문 frontmatter에 resonance를 주입·머지해 반환한다.

    _get_resonance()가 frontmatter의 `resonance:` 라인을 읽으므로,
    기록도 동일하게 frontmatter로 통일한다.
    - frontmatter 블록이 있으면: 기존 resonance 라인 교체, 없으면 닫는 --- 앞에 삽입.
    - frontmatter 블록이 없으면: 최소 frontmatter 블록을 본문 앞에 추가.
    """
    fm_match = re.match(r"^---\n(.*?)\n---\n?", content, re.DOTALL)
    if fm_match:
        fm_body = fm_match.group(1)
        if re.search(r"^resonance:\s*\S+", fm_body, re.MULTILINE):
            new_fm_body = re.sub(
                r"^resonance:\s*\S+.*$",
                f"resonance: {resonance}",
                fm_body,
                count=1,
                flags=re.MULTILINE,
            )
        else:
            new_fm_body = f"{fm_body}\nresonance: {resonance}"
        return f"---\n{new_fm_body}\n---\n" + content[fm_match.end():]

    return f"---\nresonance: {resonance}\n---\n\n{content}"


def find_unprocessed(priority_only: bool = False) -> list[Path]:
    """
    미처리 raw/ 파일 목록을 반환한다.

    priority_only=True이면 frontmatter에 resonance: high 인 파일만 반환한다.
    """
    state = load_state()
    processed = set(state.get("processed", []))
    files = [
        f for f in sorted(RAW_DIR.rglob("*"))
        if f.is_file()
        and f.suffix.lower() in SUPPORTED_EXTENSIONS
        and str(f.relative_to(WIKI_ROOT)) not in processed
    ]
    if priority_only:
        files = [f for f in files if _get_resonance(f) == "high"]
    return files


def is_duplicate(file: Path) -> bool:
    """
    index.md의 [[wikilink]] 목록과 파일명 slug를 비교해
    이미 wiki에 존재하는 주제인지 확인한다.
    중복이면 경고 메시지를 출력하고 True를 반환한다 (ingest를 중단하지는 않는다).
    """
    index_file = WIKI_ROOT / "index.md"
    if not index_file.exists():
        return False

    index_text = index_file.read_text(errors="replace")
    existing_slugs = set(re.findall(r"\[\[([^\]|/]+?)(?:\|[^\]]*)?\]\]", index_text))

    # 파일명에서 날짜 접두사(YYYY-MM-DD-) 제거 후 slug 추출
    stem = file.stem
    stem = re.sub(r"^\d{4}-\d{2}-\d{2}-", "", stem)  # 날짜 접두사 제거
    stem = stem.replace("_", "-").lower()

    if stem in {s.lower() for s in existing_slugs}:
        print(f"  [경고] '{stem}' 주제가 index.md에 이미 존재합니다.")
        print(f"         기존 wiki 페이지에 병합하는 것을 권장합니다.")
        return True
    return False


def scrape_url(url: str, resonance: str | None = None) -> Path:
    print(f"  스크랩 중: {url}")
    resp = httpx.get(url, follow_redirects=True, timeout=30,
                     headers={"User-Agent": "Mozilla/5.0"})
    resp.raise_for_status()

    md_content = markdownify(resp.text, heading_style="ATX")
    slug = re.sub(r"[^a-z0-9]+", "-", url.split("//")[-1].lower())[:60]
    date_str = datetime.now().strftime("%Y-%m-%d")
    out_file = RAW_DIR / "clippings" / f"{date_str}-{slug}.md"
    out_file.parent.mkdir(parents=True, exist_ok=True)
    resonance_line = f"resonance: {resonance}\n" if resonance else ""
    out_file.write_text(
        f"---\ntitle: 웹 스크랩\nurl: {url}\ncollected: {date_str}\n{resonance_line}---\n\n{md_content}"
    )
    print(f"  저장: {out_file.relative_to(WIKI_ROOT)}")
    return out_file


def ingest_file(src: Path, resonance: str | None = None) -> Path:
    """로컬 파일을 raw/docs/에 복사하고 텍스트 추출 MD를 함께 저장한다."""
    src = src.expanduser().resolve()
    if not src.exists():
        print(f"  오류: 파일을 찾을 수 없음 — {src}")
        sys.exit(1)
    if src.suffix.lower() not in SUPPORTED_EXTENSIONS:
        print(f"  오류: 지원하지 않는 형식 — {src.suffix}")
        sys.exit(1)

    date_str = datetime.now().strftime("%Y-%m-%d")
    docs_dir = RAW_DIR / "docs"
    docs_dir.mkdir(parents=True, exist_ok=True)

    # 원본 파일 복사
    dst = docs_dir / f"{date_str}-{src.name}"
    is_md_txt = src.suffix.lower() in {".md", ".txt"}
    if is_md_txt and resonance:
        # md/txt는 사이드카가 없으므로 복사본 frontmatter에 resonance를 주입한다.
        # (비-md/txt는 아래 .extracted.md 사이드카에 기록.)
        merged = _merge_resonance_frontmatter(src.read_text(errors="replace"), resonance)
        dst.write_text(merged)
    else:
        shutil.copy2(src, dst)

    # MD·TXT가 아닌 경우 텍스트 추출 MD도 저장
    if not is_md_txt:
        text = extract_text(src)
        if text:
            resonance_line = f"resonance: {resonance}\n" if resonance else ""
            md_out = docs_dir / f"{date_str}-{src.stem}.extracted.md"
            md_out.write_text(
                f"---\ntitle: {src.name} 추출본\nsource_file: {src.name}\nextracted: {date_str}\n{resonance_line}---\n\n{text}"
            )
            print(f"  추출 MD: {md_out.relative_to(WIKI_ROOT)}")

    print(f"  저장: {dst.relative_to(WIKI_ROOT)}")
    return dst


def save_note(text: str, resonance: str | None = None) -> Path:
    date_str = datetime.now().strftime("%Y-%m-%d-%H%M")
    out_file = RAW_DIR / "notes" / f"{date_str}-note.md"
    out_file.parent.mkdir(exist_ok=True)
    resonance_line = f"resonance: {resonance}\n" if resonance else ""
    out_file.write_text(
        f"---\ntitle: 수동 노트\ncreated: {date_str}\n{resonance_line}---\n\n{text}"
    )
    print(f"  저장: {out_file.relative_to(WIKI_ROOT)}")
    return out_file


def mark_done() -> None:
    state = load_state()
    all_files = [
        str(f.relative_to(WIKI_ROOT))
        for f in RAW_DIR.rglob("*")
        if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS
    ]
    state["processed"] = all_files
    save_state(state)
    print(f"[ingest] {len(all_files)}개 파일 처리 완료로 표시.")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--url", help="스크랩할 URL")
    parser.add_argument("--file", help="raw/docs/에 추가할 로컬 파일 경로")
    parser.add_argument("--note", help="저장할 텍스트 노트")
    parser.add_argument("--mark-done", action="store_true",
                        help="현재 raw/ 전체를 처리 완료로 표시")
    parser.add_argument(
        "--resonance",
        choices=["high", "medium", "low"],
        help="중요도 레벨 (--url/--file/--note와 함께 사용). frontmatter에 기록됨.",
    )
    parser.add_argument(
        "--priority-only",
        action="store_true",
        help="미처리 파일 중 resonance: high 파일만 출력",
    )
    args = parser.parse_args()

    if args.mark_done:
        mark_done()
        return

    if args.url:
        saved = scrape_url(args.url, resonance=args.resonance)
        is_duplicate(saved)
    elif args.file:
        saved = ingest_file(Path(args.file), resonance=args.resonance)
        is_duplicate(saved)
    elif args.note:
        saved = save_note(args.note, resonance=args.resonance)
        is_duplicate(saved)

    # 미처리 파일 목록 출력
    pending = find_unprocessed(priority_only=args.priority_only)
    if not pending:
        label = "우선순위(high) " if args.priority_only else ""
        print(f"[ingest] 처리할 새 {label}파일 없음.")
        sys.exit(0)

    label = "우선순위(high) " if args.priority_only else ""
    print(f"[ingest] 미처리 {label}파일 {len(pending)}개:")
    for f in pending:
        resonance = _get_resonance(f)
        resonance_tag = f" [{resonance}]" if resonance else ""
        print(f"  - {f.relative_to(WIKI_ROOT)}{resonance_tag}")

    # exit code 1 = 처리할 파일 있음 (run_daily.sh가 이를 감지해 LLM 호출)
    sys.exit(1)


# ── Graph delta pipeline ────────────────────────────────────

_GRAPH_FILE      = WIKI_ROOT / "wiki" / "graph.json"
_GRAPH_PREV_FILE = WIKI_ROOT / "wiki" / ".graph_prev.json"
_CANVAS_DIR      = WIKI_ROOT / "wiki" / "canvas"
_EXPORT_SCRIPT   = Path(__file__).parent / "export_graph.py"


def snapshot_graph(wiki_dir: Path | None = None) -> None:
    """wiki/graph.json → wiki/.graph_prev.json 복사. graph.json 없으면 무시."""
    _wiki = Path(wiki_dir) if wiki_dir else WIKI_ROOT / "wiki"
    src = _wiki / "graph.json"
    dst = _wiki / ".graph_prev.json"
    if src.exists():
        shutil.copy2(src, dst)


def run_delta_pipeline(wiki_dir: Path | None = None) -> dict | None:
    """
    현재 graph.json과 .graph_prev.json을 비교해 delta dict를 반환한다.
    delta가 없거나 graph.json이 없으면 None 반환.
    """
    sys.path.insert(0, str(Path(__file__).parent))
    from canvas_utils import compute_delta  # noqa: PLC0415

    _wiki = Path(wiki_dir) if wiki_dir else WIKI_ROOT / "wiki"
    cur_path  = _wiki / "graph.json"
    prev_path = _wiki / ".graph_prev.json"

    if not cur_path.exists():
        return None

    current = json.loads(cur_path.read_text())
    prev = json.loads(prev_path.read_text()) if prev_path.exists() else {"nodes": [], "links": []}

    delta = compute_delta(current, prev)
    has_changes = any([
        delta["new_nodes"], delta["removed_nodes"],
        delta["updated_nodes"], delta["new_edges"],
    ])
    return delta if has_changes else None


def print_delta(delta: dict) -> None:
    """delta를 터미널에 출력한다."""
    n_new = len(delta["new_nodes"])
    n_upd = len(delta["updated_nodes"])
    n_rem = len(delta["removed_nodes"])
    print(f"[ingest] delta — {n_new}개 신규, {n_upd}개 갱신, {n_rem}개 제거")

    for node in delta["new_nodes"]:
        cat = node.get("category") or "?"
        print(f"  + {node['id']}  ({cat}/)  inbound 0 → {node['inbound']}")

    for node in delta["updated_nodes"]:
        cat = node.get("category") or "?"
        print(f"  ~ {node['id']}  ({cat}/)  inbound → {node['inbound']}")

    for node in delta["removed_nodes"]:
        cat = node.get("category") or "?"
        print(f"  - {node['id']}  ({cat}/)  제거됨")

    new_edges = delta["new_edges"]
    if new_edges:
        first = new_edges[0]
        rest  = len(new_edges) - 1
        msg = f"  엣지 +{len(new_edges)}: {first['source']} → {first['target']}"
        if rest > 0:
            msg += f" 외 {rest}개"
        print(msg)


def generate_ingest_delta_canvas(wiki_dir: Path | None = None) -> bool:
    """
    delta canvas를 wiki/canvas/ingest-delta.canvas로 저장한다.
    snapshot_graph() + export_graph.py 실행 이후에 호출해야 한다.
    canvas가 생성되면 True, delta 없으면 False 반환.
    """
    sys.path.insert(0, str(Path(__file__).parent))
    from canvas_utils import build_delta_canvas, save_canvas  # noqa: PLC0415

    _wiki = Path(wiki_dir) if wiki_dir else WIKI_ROOT / "wiki"
    cur_path  = _wiki / "graph.json"
    prev_path = _wiki / ".graph_prev.json"

    if not cur_path.exists():
        return False

    current = json.loads(cur_path.read_text())
    prev = json.loads(prev_path.read_text()) if prev_path.exists() else {"nodes": [], "links": []}

    canvas = build_delta_canvas(current, prev)
    if canvas is None:
        return False

    out_path = _wiki / "canvas" / "ingest-delta.canvas"
    save_canvas(canvas, out_path)
    print(f"[ingest] canvas → wiki/canvas/ingest-delta.canvas")
    return True


if __name__ == "__main__":
    main()
